"""Build the enhanced C242 final-presentation deck.

Strategy
--------
- Slide size 720 x 405 pt (16:9, identical to the source deck).
- Every original page from `docs/presentations/C242_Final.pdf` is preserved verbatim by embedding
  the corresponding rendered PNG full-bleed on its own slide.
- Supplementary slides are interleaved (and appended) to fill known gaps:
  cover, agenda, motivation, dataset body, preprocessing pipeline, training
  setup, loss curves, PR curve, frequency analysis, per-ontology summaries,
  team progression, conclusion, future work, Q&A.
- Visual consistency: white background, navy titles (#1F3A68), light grey
  bottom rule, charcoal body text, restrained accent palette.

Run with the cafa5 conda env:
    /global/home/users/bensonli/.conda/envs/cafa5/bin/python scripts/build_enhanced_deck.py
"""

from __future__ import annotations

import os
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Pt, Emu
from PIL import Image


REPO = Path("/global/home/users/bensonli/c242_cafa5")
PAGES_DIR = Path("/tmp/c242_pages")
FIGS = REPO / "figures"
OUT = REPO / "docs/presentations/C242_Final_enhanced.pptx"

NAVY = RGBColor(0x1F, 0x3A, 0x68)
ACCENT = RGBColor(0xE0, 0x7A, 0x1F)        # warm orange (matches bar charts)
ACCENT_BLUE = RGBColor(0x3B, 0x77, 0xC4)
CHARCOAL = RGBColor(0x33, 0x33, 0x33)
GREY = RGBColor(0x88, 0x8A, 0x8E)
LIGHT_GREY = RGBColor(0xE2, 0xE5, 0xE9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Pt(720)
SLIDE_H = Pt(405)


def _add_textbox(slide, left, top, width, height, text, *,
                 size=18, bold=False, color=CHARCOAL,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                 font="Calibri"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    tf.vertical_anchor = anchor
    lines = text if isinstance(text, list) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if isinstance(line, tuple):
            text_value, opts = line
        else:
            text_value, opts = line, {}
        run = p.add_run()
        run.text = text_value
        f = run.font
        f.name = font
        f.size = Pt(opts.get("size", size))
        f.bold = opts.get("bold", bold)
        c = opts.get("color", color)
        f.color.rgb = c
    return box


def _add_bullets(slide, left, top, width, height, bullets,
                 *, size=16, color=CHARCOAL, line_spacing=1.15):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        if isinstance(b, tuple):
            txt, opts = b
        else:
            txt, opts = b, {}
        level = opts.get("level", 0)
        bullet_char = "•" if level == 0 else "–"
        run = p.add_run()
        run.text = f"{'  ' * level}{bullet_char}  {txt}"
        f = run.font
        f.name = "Calibri"
        f.size = Pt(opts.get("size", size if level == 0 else max(12, size - 2)))
        f.bold = opts.get("bold", False)
        f.color.rgb = opts.get("color", color)
    return box


def _add_rect(slide, left, top, width, height, fill=NAVY, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def _header(slide, title, subtitle=None):
    """Standard header for supplementary slides."""
    # Top accent rule
    _add_rect(slide, Pt(0), Pt(0), SLIDE_W, Pt(6), fill=NAVY)
    # Title
    _add_textbox(slide, Pt(36), Pt(18), Pt(648), Pt(34),
                 title, size=24, bold=True, color=NAVY)
    if subtitle:
        _add_textbox(slide, Pt(36), Pt(50), Pt(648), Pt(20),
                     subtitle, size=13, color=GREY, font="Calibri")
    # Bottom accent line under title
    _add_rect(slide, Pt(36), Pt(74), Pt(72), Pt(2), fill=ACCENT)


def _footer(slide, page_label):
    _add_textbox(slide, Pt(36), Pt(388), Pt(400), Pt(14),
                 "C242 Final  ·  CAFA-5 Protein Function Prediction",
                 size=9, color=GREY)
    _add_textbox(slide, Pt(640), Pt(388), Pt(48), Pt(14),
                 page_label, size=9, color=GREY, align=PP_ALIGN.RIGHT)


def _add_picture_fit(slide, image_path, left, top, max_w, max_h, *,
                     center=True):
    """Embed an image scaled to fit in the given box, preserving aspect ratio."""
    with Image.open(image_path) as im:
        iw, ih = im.size
    box_w_pt = max_w / 12700.0  # EMU → pt? actually pptx uses EMU under hood
    # python-pptx accepts Pt() or Emu(); use the EMU values directly.
    scale = min(max_w / iw, max_h / ih)
    w = int(iw * scale)
    h = int(ih * scale)
    if center:
        left_eff = left + (max_w - w) // 2
        top_eff = top + (max_h - h) // 2
    else:
        left_eff, top_eff = left, top
    return slide.shapes.add_picture(str(image_path), left_eff, top_eff,
                                    width=w, height=h)


def _blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank layout


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def slide_cover(prs):
    s = _blank_slide(prs)
    # Full navy band on left
    _add_rect(s, Pt(0), Pt(0), Pt(260), SLIDE_H, fill=NAVY)
    _add_rect(s, Pt(260), Pt(0), Pt(6), SLIDE_H, fill=ACCENT)

    _add_textbox(s, Pt(28), Pt(60), Pt(220), Pt(20),
                 "C242  ·  Final Presentation",
                 size=12, color=WHITE)
    _add_textbox(s, Pt(28), Pt(90), Pt(220), Pt(20),
                 "Spring 2026", size=11, color=LIGHT_GREY)

    _add_textbox(s, Pt(290), Pt(110), Pt(400), Pt(60),
                 "Can structural information improve",
                 size=28, bold=True, color=NAVY)
    _add_textbox(s, Pt(290), Pt(146), Pt(400), Pt(40),
                 "protein function prediction beyond",
                 size=28, bold=True, color=NAVY)
    _add_textbox(s, Pt(290), Pt(180), Pt(400), Pt(40),
                 "sequence-only features?",
                 size=28, bold=True, color=ACCENT)

    _add_textbox(s, Pt(290), Pt(238), Pt(400), Pt(20),
                 "A controlled comparison on the CAFA-5 challenge",
                 size=14, color=CHARCOAL)
    _add_textbox(s, Pt(290), Pt(258), Pt(400), Pt(20),
                 "Sequence baseline · AlphaFold structure graphs · Late fusion",
                 size=12, color=GREY)

    _add_rect(s, Pt(290), Pt(298), Pt(48), Pt(2), fill=ACCENT)
    _add_textbox(s, Pt(290), Pt(308), Pt(400), Pt(20),
                 "Team C242 — CAFA-5 Project",
                 size=12, bold=True, color=NAVY)
    _add_textbox(s, Pt(290), Pt(326), Pt(400), Pt(18),
                 "ESM2 protein-LM  +  GCN structure GNN  +  late fusion classifier",
                 size=10, color=GREY)


def slide_agenda(prs):
    s = _blank_slide(prs)
    _header(s, "Agenda", "What we will walk through today")
    items = [
        "1.  Motivation & research question",
        "2.  CAFA-5 task and data acquisition",
        "3.  Three challenges in the dataset",
        "4.  AlphaFold structure pipeline",
        "5.  Graph features  (node / edge / graph)",
        "6.  Dataset splits and label vocabulary",
        "7.  Model architecture and training setup",
        "8.  Results — controlled MFO comparison",
        "9.  Per-ontology and ablation results",
        "10.  Conclusions and future work",
    ]
    _add_textbox(s, Pt(60), Pt(105), Pt(600), Pt(280),
                 [(t, {"size": 18, "color": CHARCOAL}) for t in items],
                 size=18)
    _footer(s, "")


def slide_motivation(prs):
    s = _blank_slide(prs)
    _header(s, "Motivation",
            "Why ask about structure when sequence models already work?")
    _add_bullets(s, Pt(36), Pt(95), Pt(380), Pt(280), [
        ("Protein sequence ⇒ function is one of biology's central problems.",
         {"size": 16, "bold": True, "color": NAVY}),
        ("Sequence-based protein language models (ESM2) already give a strong baseline on CAFA-5.",
         {"size": 14}),
        ("BUT function is ultimately determined by 3-D structure: active sites, "
         "binding pockets, and folds.",
         {"size": 14}),
        ("AlphaFold-2 / AF-DB now makes high-quality predicted structures "
         "available for nearly every UniProt entry.",
         {"size": 14}),
        ("Question: does structural information add a measurable signal beyond "
         "sequence-only features on the CAFA-5 benchmark?",
         {"size": 14, "bold": True, "color": ACCENT}),
    ], size=14)

    # Right: highlight box
    _add_rect(s, Pt(440), Pt(110), Pt(244), Pt(220), fill=RGBColor(0xF4, 0xF6, 0xFA))
    _add_rect(s, Pt(440), Pt(110), Pt(4), Pt(220), fill=ACCENT)
    _add_textbox(s, Pt(454), Pt(122), Pt(220), Pt(20),
                 "Research question",
                 size=11, bold=True, color=ACCENT)
    _add_textbox(s, Pt(454), Pt(140), Pt(220), Pt(80),
                 "Whether structural information can improve "
                 "function prediction beyond sequence-only features?",
                 size=15, bold=True, color=NAVY)
    _add_textbox(s, Pt(454), Pt(220), Pt(220), Pt(20),
                 "Approach",
                 size=11, bold=True, color=ACCENT)
    _add_textbox(s, Pt(454), Pt(238), Pt(220), Pt(80),
                 "Build a sequence-only baseline and a structure-graph "
                 "model on the same protein cohort, then compare and fuse them.",
                 size=12, color=CHARCOAL)
    _footer(s, "")


def slide_dataset_body(prs):
    """Fills in the empty 'Data acquisition' slide (page 3) with detail."""
    s = _blank_slide(prs)
    _header(s, "Data acquisition",
            "CAFA-5 protein-function-prediction challenge")
    _add_bullets(s, Pt(36), Pt(95), Pt(380), Pt(280), [
        ("CAFA-5 (Kaggle, 2023) — predict GO terms for ~140 K UniProt proteins.",
         {"size": 14, "bold": True}),
        ("Inputs: train_sequences.fasta, train_terms.tsv (sequence ⇒ GO list), "
         "test_superset.fasta, taxonomy info.",
         {"size": 13}),
        ("Three sub-ontologies: MFO (Molecular Function), CCO (Cellular "
         "Component), BPO (Biological Process).",
         {"size": 13}),
        ("Structure source: AlphaFold-2 predictions (PDB + PAE) downloaded "
         "from EBI AF-DB, matched to CAFA proteins by UniProt accession.",
         {"size": 13}),
        ("Output: matched training_index.parquet linking each CAFA protein to "
         "its AlphaFold artifacts and fragment list.",
         {"size": 13}),
    ], size=13)

    # Stat callouts
    stat_specs = [
        ("≈ 140 K", "CAFA-5 proteins"),
        ("≈ 35 K", "Unique GO terms"),
        ("≈ 5.4 M", "Annotation pairs"),
        ("≈ 90%", "AlphaFold match"),
    ]
    x = Pt(440)
    y = Pt(110)
    cell_w = Pt(120)
    cell_h = Pt(70)
    for i, (big, small) in enumerate(stat_specs):
        col = i % 2
        row = i // 2
        cx = x + col * (cell_w + Pt(6))
        cy = y + row * (cell_h + Pt(8))
        _add_rect(s, cx, cy, cell_w, cell_h, fill=RGBColor(0xF4, 0xF6, 0xFA))
        _add_rect(s, cx, cy, Pt(3), cell_h, fill=ACCENT)
        _add_textbox(s, cx + Pt(10), cy + Pt(8), cell_w - Pt(12), Pt(28),
                     big, size=22, bold=True, color=NAVY)
        _add_textbox(s, cx + Pt(10), cy + Pt(40), cell_w - Pt(12), Pt(20),
                     small, size=10, color=GREY)
    _add_textbox(s, Pt(440), Pt(280), Pt(244), Pt(40),
                 "Numbers are approximate after AlphaFold matching and "
                 "deduplication on the matched cohort.",
                 size=9, color=GREY)
    _footer(s, "")


def slide_pipeline(prs):
    s = _blank_slide(prs)
    _header(s, "End-to-end pipeline",
            "From raw CAFA-5 inputs to evaluated GO predictions")
    _add_picture_fit(s, FIGS / "preprocessing_pipeline_diagram.png",
                     Pt(36), Pt(96), Pt(648), Pt(220))
    _add_bullets(s, Pt(36), Pt(322), Pt(648), Pt(60), [
        ("Each stage is a deterministic, cached artifact — graph caches "
         "support reproducible re-training without recomputing features.",
         {"size": 11, "color": GREY}),
    ], size=11)
    _footer(s, "")


def slide_training_setup(prs):
    s = _blank_slide(prs)
    _header(s, "Training setup",
            "Identical baseline + tuned configuration (Savio 1080Ti)")
    _add_picture_fit(s, FIGS / "training_setup_table.png",
                     Pt(36), Pt(92), Pt(648), Pt(258))
    _add_textbox(s, Pt(36), Pt(360), Pt(648), Pt(20),
                 "Tuned column lifts CCO Fmax by +2.28 pp via wider hidden "
                 "size, weighted-BCE, and label-aware scoring.",
                 size=11, color=GREY)
    _footer(s, "")


def slide_loss_curves(prs):
    s = _blank_slide(prs)
    _header(s, "Validation curves by epoch",
            "Raw graph baseline vs. label-aware scorer (longer confirmation run)")
    _add_picture_fit(s, FIGS / "training_validation_loss_curves.png",
                     Pt(36), Pt(92), Pt(648), Pt(240))
    _add_bullets(s, Pt(36), Pt(338), Pt(648), Pt(50), [
        ("CCO: tuned label-aware run reaches 0.586 by epoch 6 (+2.4 pp over baseline).",
         {"size": 12, "bold": True}),
        ("MFO: gains are tighter (+0.1 pp) — diminishing returns dominate the "
         "MFO label distribution.", {"size": 12}),
    ], size=12)
    _footer(s, "")


def slide_pr_curve(prs):
    s = _blank_slide(prs)
    _header(s, "MFO precision–recall (micro)",
            "Where exactly does structural information help?")
    _add_picture_fit(s, FIGS / "mfo_precision_recall_curve.png",
                     Pt(36), Pt(92), Pt(420), Pt(260))
    _add_bullets(s, Pt(470), Pt(110), Pt(214), Pt(260), [
        ("AUPRC: 0.359 → 0.367 (+0.80 pp)",
         {"size": 13, "bold": True, "color": ACCENT}),
        ("Largest precision lift (≈ +4 pp) lands in the recall = 0.35–0.45 "
         "operating range — exactly where CAFA-5 evaluates Fmax.",
         {"size": 11}),
        ("Curves coincide at very low and very high recall, which matches "
         "intuition: structure cannot help when no signal exists, and head "
         "labels are already saturated by ESM2.",
         {"size": 11}),
    ], size=11)
    _footer(s, "")


def slide_freq_analysis(prs):
    s = _blank_slide(prs)
    _header(s, "Where does the gain come from?",
            "Per-term F1 by training-set GO frequency")
    _add_picture_fit(s, FIGS / "performance_vs_go_term_frequency.png",
                     Pt(36), Pt(92), Pt(648), Pt(240))
    _add_bullets(s, Pt(36), Pt(338), Pt(648), Pt(50), [
        ("Gain (+0.90 pp) concentrates on high-frequency GO terms (>1000 examples).",
         {"size": 12, "bold": True}),
        ("Long-tail terms remain effectively zero for both models — a "
         "label-imbalance limit that motivates weighted-BCE and label-aware "
         "scoring as next steps.", {"size": 12}),
    ], size=12)
    _footer(s, "")


def slide_team_baselines(prs):
    s = _blank_slide(prs)
    _header(s, "Sequence-baseline progression",
            "From classical k-mer features to ESM2 protein-LM embeddings")
    _add_picture_fit(s, FIGS / "team_progress_sequence_baselines.png",
                     Pt(36), Pt(92), Pt(648), Pt(240))
    _add_bullets(s, Pt(36), Pt(338), Pt(648), Pt(50), [
        ("ESM2 lifts micro-F1 by +13.3 pp on MFO, +11.7 pp on CCO, +8.4 pp on BPO.",
         {"size": 12, "bold": True}),
        ("This is the strong ESM2 baseline against which the structure-graph "
         "model must compete.", {"size": 12}),
    ], size=12)
    _footer(s, "")


def slide_per_ontology_combined(prs):
    s = _blank_slide(prs)
    _header(s, "Per-ontology graph experiments",
            "Raw graph  ·  Weighted BCE  ·  Label-aware scorer")
    _add_picture_fit(s, FIGS / "per_ontology_graph_comparison.png",
                     Pt(36), Pt(92), Pt(648), Pt(240))
    _add_bullets(s, Pt(36), Pt(338), Pt(648), Pt(50), [
        ("CCO benefits most from label-aware scoring (+2.28 pp) — fewer, "
         "less-skewed labels.",
         {"size": 12, "bold": True}),
        ("MFO benefits most from weighted-BCE (+0.31 pp); label-aware "
         "scoring barely moves the needle.",
         {"size": 12}),
    ], size=12)
    _footer(s, "")


def slide_node_edge_table(prs):
    s = _blank_slide(prs)
    _header(s, "Graph feature inventory",
            "What lives at residues, contacts, and the whole protein")
    _add_picture_fit(s, FIGS / "node_edge_graph_data_comparison.png",
                     Pt(36), Pt(96), Pt(648), Pt(260))
    _footer(s, "")


def slide_conclusion(prs):
    s = _blank_slide(prs)
    _header(s, "Conclusion",
            "Answering the research question with the controlled experiment")
    _add_bullets(s, Pt(36), Pt(95), Pt(648), Pt(260), [
        ("Yes — structural information improves CAFA-5 prediction beyond "
         "sequence features, but the absolute gain is small.",
         {"size": 16, "bold": True, "color": ACCENT}),
        ("Matched MFO test Fmax: 0.4557 (sequence-only) → 0.4580 (graph) → "
         "0.4579 (late fusion).",
         {"size": 13}),
        ("Micro-AUPRC: 0.359 → 0.367 (+0.80 pp); the precision lift "
         "concentrates in the operating range CAFA-5 actually scores.",
         {"size": 13}),
        ("Per-ontology, the largest absolute structural gains appear on CCO "
         "with label-aware scoring (+2.28 pp Fmax).",
         {"size": 13}),
        ("Late fusion with α tuned on validation does not beat the structure "
         "branch alone — sequence features are already saturating the gradient.",
         {"size": 13}),
        ("Overall: structure helps where the label has enough support; "
         "long-tail GO terms are bottlenecked by data, not modality.",
         {"size": 13, "bold": True}),
    ], size=13)
    _footer(s, "")


def slide_future_work(prs):
    s = _blank_slide(prs)
    _header(s, "Limitations & future work",
            "What we would change with another month")
    _add_bullets(s, Pt(36), Pt(95), Pt(648), Pt(260), [
        ("Long-tail terms (≤ 20 examples) are essentially zero — try "
         "hierarchical loss that propagates evidence up the GO DAG.",
         {"size": 13, "bold": True}),
        ("Replace 2-layer GCN with an SE(3)-equivariant or distance-aware GNN "
         "(GVP, SchNet) to better exploit Cα geometry.", {"size": 13}),
        ("Use ESM-Fold / ESM-IF embeddings as node features so the structure "
         "branch starts from a stronger inductive base.", {"size": 13}),
        ("Train on the full (≈140 K) cohort with mixed-precision and "
         "gradient-checkpointing — the current matched cohort is "
         "AlphaFold-coverage-limited.", {"size": 13}),
        ("Late fusion: replace fixed α with a learned per-label gate; current "
         "global α saturates at the structure branch.", {"size": 13}),
        ("Submit predictions to the live CAFA-5 leaderboard for an "
         "independent generalization check.", {"size": 13}),
    ], size=13)
    _footer(s, "")


def slide_qa(prs):
    s = _blank_slide(prs)
    _add_rect(s, Pt(0), Pt(0), SLIDE_W, SLIDE_H, fill=NAVY)
    _add_rect(s, Pt(0), Pt(190), SLIDE_W, Pt(4), fill=ACCENT)
    _add_textbox(s, Pt(0), Pt(140), SLIDE_W, Pt(50),
                 "Thank you  ·  Questions?",
                 size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _add_textbox(s, Pt(0), Pt(204), SLIDE_W, Pt(20),
                 "C242 Final  ·  CAFA-5 Protein Function Prediction",
                 size=14, color=LIGHT_GREY, align=PP_ALIGN.CENTER)
    _add_textbox(s, Pt(0), Pt(228), SLIDE_W, Pt(20),
                 "Code & artifacts: c242_cafa5 repository",
                 size=11, color=LIGHT_GREY, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Original-page embedding
# ---------------------------------------------------------------------------

def embed_original(prs, page_idx):
    """Embed an original PDF page full-bleed so nothing is lost."""
    img_path = PAGES_DIR / f"page_{page_idx:02d}.png"
    s = _blank_slide(prs)
    # Background fill (white) then the page image scaled to fill.
    pic = s.shapes.add_picture(str(img_path), 0, 0,
                               width=SLIDE_W, height=SLIDE_H)
    # Tag with a tiny label so reviewers know this is an original slide.
    _add_rect(s, Pt(0), Pt(395), Pt(180), Pt(10), fill=WHITE)
    _add_textbox(s, Pt(8), Pt(394), Pt(180), Pt(10),
                 f"original deck · page {page_idx}",
                 size=8, color=GREY)
    return s


# ---------------------------------------------------------------------------
# Deck assembly
# ---------------------------------------------------------------------------

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 1.  Cover + agenda
    slide_cover(prs)
    slide_agenda(prs)

    # 2.  Motivation precedes the original "research question" page.
    slide_motivation(prs)
    embed_original(prs, 1)   # blank title page (kept as-is for completeness)
    embed_original(prs, 2)   # research-question slide

    # 3.  Data acquisition: original empty-body slide + filled-in supplement.
    embed_original(prs, 3)
    slide_dataset_body(prs)
    slide_pipeline(prs)

    # 4.  Three challenges (original deck pages 4-6) with sequence baseline
    #     positioned between to motivate "why we need more than sequence".
    embed_original(prs, 4)
    embed_original(prs, 5)
    embed_original(prs, 6)
    slide_team_baselines(prs)

    # 5.  AlphaFold matching + feature breakdowns (original pages 7-10),
    #     followed by the consolidated node/edge/graph table.
    embed_original(prs, 7)
    embed_original(prs, 8)
    embed_original(prs, 9)
    embed_original(prs, 10)
    embed_original(prs, 11)            # original node/edge/graph table
    slide_node_edge_table(prs)         # cleaner branded version

    # 6.  Splits + model architecture
    embed_original(prs, 12)
    embed_original(prs, 13)
    slide_training_setup(prs)
    slide_loss_curves(prs)

    # 7.  Results: original matched MFO + supplemental PR curve and frequency.
    embed_original(prs, 14)
    slide_pr_curve(prs)
    slide_freq_analysis(prs)

    # 8.  Per-ontology results.
    embed_original(prs, 15)
    slide_per_ontology_combined(prs)

    # 9.  Conclusion: original empty conclusion page + filled-in supplement.
    embed_original(prs, 16)
    slide_conclusion(prs)
    slide_future_work(prs)
    slide_qa(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"saved → {OUT}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
